import logging
import os
import pptx
from aiogram.dispatcher import Dispatcher, FSMContext
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, InputMediaPhoto, ReplyKeyboardMarkup, \
    KeyboardButton, ReplyKeyboardRemove
from aiogram.utils.executor import start_polling
from aiogram import Bot, types

from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher.filters.state import StatesGroup, State

storage = MemoryStorage()
bot = Bot(token=os.getenv('BOT_TOKEN'))
dp = Dispatcher(bot, storage=storage)

# HEROKU_APP_NAME = os.getenv('HEROKU_APP_NAME')

# # webhook settings
# WEBHOOK_HOST = f'https://{HEROKU_APP_NAME}.herokuapp.com'
# WEBHOOK_PATH = f'/webhook/{TOKEN}'
# WEBHOOK_URL = f'{WEBHOOK_HOST}{WEBHOOK_PATH}'

# # webserver settings
# WEBAPP_HOST = '0.0.0.0'
# WEBAPP_PORT = os.getenv('PORT', default=8000)

MESSAGES = {
    "start": ["Привет, я бот, помогающий сделать презентацию в привычном формате .pptx прямо в Telegram!",
              "Введите /help, чтобы посмотреть список доступных команд"],
    "help": ["Вот список доступных команд:", "/new – Создать новую презентацию"],
    "presentation_start": "Введите название презентации (так будет называться .pptx файл)",
    "get_name": "Отлично! Ваш файл будет называться {}.pptx",
    "get_design_type": "Теперь нужно выбрать подходящий дизайн:",
    "get_slides_number": "Сколько слайдов должно быть в презентации?",
    "start_slides": "Итак, начнём заполнять слайды",
    "make_types": "Выберите тип слайда №{}",
    "make_slide": "Введите текст №{}",
    "make_pptx": "Принято! Отправил презентацию на кухню:)",
    "cancel": "Создание презентации отменено",
    "cancel_button": u'\U0000274C' + " Отменить создание презентации",

}

# prs = pptx.Presentation("designs/3.pptx")
#
# prs.slides.add_slide(prs.slide_layouts[8])
#
# shapes = prs.slides[0].shapes
#
# img = Image.open("designs/demo/2.png")
# img.thumbnail(size=(720, 390))
# buf = io.BytesIO()
# img.save(buf, format='PNG')
#
# image = shapes.add_picture(buf, left=0, top=0)
#
# image.left = (prs.slide_width - image.width) // 2
#
# prs.slides[0].placeholders[0].text = "gg"
#
# prs.save("GG.pptx")

ERRORS = {
    "int_required": "Ошибка. Введите натуральное число"
}

DESIGNS_NUMBER = len(os.listdir("designs")) - 1
SLIDE_TYPES = {
    0: 2,
    1: 2,
    2: 2,
    3: 3,
    4: 5,
    5: 1,
    6: 4,
    7: 3
}

TYPE_NUMBER = len(SLIDE_TYPES) - 1

INLINE_MARKUP_DESIGNS = InlineKeyboardMarkup()
INLINE_MARKUP_DESIGNS.row(InlineKeyboardButton(text="Назад", callback_data="back " + str(DESIGNS_NUMBER - 1)),
                          InlineKeyboardButton(text="Вперёд", callback_data="next 1"))
INLINE_MARKUP_DESIGNS.add(InlineKeyboardButton(text="Этот", callback_data="0"))
INLINE_MARKUP_DESIGNS.add(InlineKeyboardButton(text=MESSAGES["cancel_button"], callback_data="cancel"))

INLINE_MARKUP_TYPES = InlineKeyboardMarkup()
INLINE_MARKUP_TYPES.row(InlineKeyboardButton(text="Назад", callback_data="back " + str(TYPE_NUMBER)),
                        InlineKeyboardButton(text="Вперёд", callback_data="next 1"))
INLINE_MARKUP_TYPES.add(InlineKeyboardButton(text="Этот", callback_data="0"))
INLINE_MARKUP_TYPES.add(InlineKeyboardButton(text=MESSAGES["cancel_button"], callback_data="cancel"))

KEYBOARD_CANCEL = ReplyKeyboardMarkup(resize_keyboard=True)
KEYBOARD_CANCEL.add(KeyboardButton(MESSAGES["cancel_button"]))


def make_presentation(design, slides, naming):
    prs = pptx.Presentation(f"designs/{design}.pptx")
    for slide_info in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[slide_info["type"]])
        for i in range(SLIDE_TYPES[slide_info["type"]]):
            slide.placeholders[i].text = str(slide_info["content"][i])
    prs.save(naming)


async def check_cancel(message, state):
    """ Проверка кнопки клавиатуры на сброс """
    if message.text == MESSAGES["cancel_button"]:
        await message.answer(MESSAGES["cancel"], reply_markup=ReplyKeyboardRemove())
        await state.finish()
        return True
    return False


class PresentationStates(StatesGroup):
    name = State()
    design_type = State()
    slides_number = State()
    types_message = State()
    last_type = State()
    slides = State()


# async def on_startup(dispatcher):
#     await bot.set_webhook(WEBHOOK_URL, drop_pending_updates=True)


# async def on_shutdown(dispatcher):
#     await bot.delete_webhook()


@dp.message_handler(commands=["start"])
async def bot_start(message: types.Message):
    """ Стандартная команда начала диалога """
    for msg in MESSAGES["start"]:
        await message.answer(msg)


@dp.message_handler(commands=["help"])
async def bot_help(message: types.Message):
    """ Показ пользователю доступных команд """
    for msg in MESSAGES["help"]:
        await message.answer(msg)


@dp.message_handler(commands=["new"])
async def presentation_start(message: types.Message):
    """ Начало создания презентации – запрос названия файла """
    try:
        logging.log(logging.INFO, str(message.from_user.id) + " " + message.from_user.username)
    except Exception:
        logging.log(logging.INFO, message.from_user.id)
    await message.answer(MESSAGES["presentation_start"], reply_markup=KEYBOARD_CANCEL)
    await PresentationStates.name.set()


@dp.message_handler(state=PresentationStates.name)
async def get_filename(message: types.Message, state: FSMContext):
    """ Получение названия от пользователя """
    if await check_cancel(message, state):
        return
    await state.update_data(name=message.text)
    await message.answer(MESSAGES["get_name"].format(message.text))
    await PresentationStates.design_type.set()
    await message.answer(MESSAGES["get_design_type"], reply_markup=ReplyKeyboardRemove())
    await bot.send_photo(chat_id=message.chat.id, photo=open("designs/demo/0.png", "rb"), caption="Дизайн №1",
                         reply_markup=INLINE_MARKUP_DESIGNS)


@dp.callback_query_handler(lambda c: True, state=PresentationStates.design_type)
async def inline_designs(call: types.CallbackQuery, state: FSMContext):
    """ Логика inline-клавиатуры под фото дизайнов """
    await bot.answer_callback_query(call.id)
    if call.data == "cancel":
        await bot.edit_message_caption(caption=MESSAGES["cancel"],
                                       chat_id=call.message.chat.id,
                                       message_id=call.message.message_id,
                                       reply_markup=None)
        await state.finish()
    elif call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = InlineKeyboardMarkup()
        if call.data.split()[1] == str(DESIGNS_NUMBER - 1):
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next 0")
        else:
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next " + str(int(call.data.split()[1]) + 1))
        if call.data.split()[1] == "0":
            mark.row(InlineKeyboardButton(text='Назад', callback_data="back " + str(DESIGNS_NUMBER - 1)),
                     but)
        else:
            mark.row(InlineKeyboardButton(text='Назад',
                                          callback_data="back " + str(int(call.data.split()[1]) - 1)),
                     but)
        mark.add(InlineKeyboardButton(text='Этот', callback_data=call.data.split()[1]))
        mark.add(InlineKeyboardButton(text=MESSAGES["cancel_button"], callback_data="cancel"))
        with open(f"designs/demo/{call.data.split()[1]}.png", "rb") as file:
            await bot.edit_message_media(chat_id=call.message.chat.id,
                                         message_id=call.message.message_id,
                                         media=InputMediaPhoto(media=file,
                                                               caption=f"Дизайн №{str(int(call.data.split()[1]) + 1)}"),
                                         reply_markup=mark)
    else:
        await bot.edit_message_caption(caption=f"Дизайн вашей презентации (№{str(int(call.data) + 1)}), выбранный вами",
                                       chat_id=call.message.chat.id,
                                       message_id=call.message.message_id)
        await state.update_data(design_type=int(call.data))
        await call.message.answer(MESSAGES["get_slides_number"], reply_markup=KEYBOARD_CANCEL)
        await PresentationStates.slides_number.set()


@dp.message_handler(state=PresentationStates.slides_number)
async def get_number(message: types.Message, state: FSMContext):
    """ Получение количества слайдов от пользователя """
    if await check_cancel(message, state):
        return
    try:
        if float(message.text) != float(int(message.text)):
            raise ValueError
        if int(message.text) <= 0:
            raise ValueError
        await state.update_data(slides_number=int(message.text))
        await message.answer(MESSAGES["start_slides"])
        await message.answer(MESSAGES["make_types"].format(1), reply_markup=ReplyKeyboardRemove())
        await state.update_data(types_message=0)
        await PresentationStates.last_type.set()
        await bot.send_photo(chat_id=message.chat.id, photo=open("types/0.png", "rb"), caption="Тип №1",
                             reply_markup=INLINE_MARKUP_TYPES)
    except ValueError:
        await message.answer(ERRORS["int_required"])
        await PresentationStates.slides_number.set()


@dp.callback_query_handler(lambda c: True, state=PresentationStates.last_type)
async def inline_types(call: types.CallbackQuery, state: FSMContext):
    """ Логика inline-клавиатуры под фото типов слайдов """
    await bot.answer_callback_query(call.id)
    data = await state.get_data()
    if data["types_message"] == 0:
        await state.update_data(types_message=call.message.message_id)
    if call.data == "cancel":
        await bot.edit_message_caption(caption=MESSAGES["cancel"],
                                       chat_id=call.message.chat.id,
                                       message_id=call.message.message_id,
                                       reply_markup=None)
        await state.finish()
    elif call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = InlineKeyboardMarkup()
        if call.data.split()[1] == str(TYPE_NUMBER - 1):
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next 0")
        else:
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next " + str(int(call.data.split()[1]) + 1))
        if call.data.split()[1] == "0":
            mark.row(InlineKeyboardButton(text='Назад', callback_data="back " + str(TYPE_NUMBER - 1)),
                     but)
        else:
            mark.row(InlineKeyboardButton(text='Назад',
                                          callback_data="back " + str(int(call.data.split()[1]) - 1)),
                     but)
        mark.add(InlineKeyboardButton(text='Этот', callback_data=call.data.split()[1]))
        mark.add(InlineKeyboardButton(text=MESSAGES["cancel_button"], callback_data="cancel"))
        try:
            with open(f"types/{call.data.split()[1]}.png", "rb") as file:
                await bot.edit_message_media(chat_id=call.message.chat.id,
                                             message_id=call.message.message_id,
                                             media=InputMediaPhoto(media=file,
                                                                   caption=f"Тип №{str(int(call.data.split()[1]) + 1)}"),
                                             reply_markup=mark)
        except FileNotFoundError:
            """ Так как нет 6-ого типа слайда """
            with open(f"types/{str(int(call.data.split()[1]) + 1)}.png", "rb") as file:
                await bot.edit_message_media(chat_id=call.message.chat.id,
                                             message_id=call.message.message_id,
                                             media=InputMediaPhoto(media=file,
                                                                   caption=f"Тип №{str(int(call.data.split()[1]) + 1)}"),
                                             reply_markup=mark)
    else:
        try:
            slides_len = len(data["slides"])
        except KeyError:
            slides_len = 0
        slides_len += 1
        await bot.edit_message_caption(
            caption=f"Тип слайда (№{str(int(call.data) + 1)}), выбранный вами для слайда №{slides_len}",
            chat_id=call.message.chat.id,
            message_id=call.message.message_id)
        await call.message.answer(MESSAGES["make_slide"].format(1), reply_markup=KEYBOARD_CANCEL)
        await state.update_data(last_type=int(call.data))
        await PresentationStates.slides.set()


@dp.message_handler(state=PresentationStates.slides)
async def make_slides(message: types.Message, state: FSMContext):
    if await check_cancel(message, state):
        return
    data = await state.get_data()
    edited = False
    try:
        slides = data["slides"]
        last_slide = slides[-1]
        if SLIDE_TYPES[last_slide["type"]] == len(last_slide["content"]):
            current_slide = {}
        else:
            current_slide = last_slide
            edited = True
    except KeyError:
        slides = []
        current_slide = {}
    current_slide["type"] = data["last_type"]

    if edited:
        current_slide["content"].append(message.text)
        slides[-1] = current_slide
    else:
        current_slide["content"] = [message.text]
        slides.append(current_slide)

    await state.update_data(slides=slides)

    needed_number = SLIDE_TYPES[slides[-1]["type"]]
    now_number = len(slides[-1]["content"])
    if needed_number == now_number:
        if len(slides) == data["slides_number"]:
            await message.answer(MESSAGES["make_pptx"], reply_markup=ReplyKeyboardRemove())
            make_presentation(data["design_type"], slides, data["name"] + '.pptx')
            await bot.send_document(chat_id=message.chat.id, document=open(data["name"] + '.pptx', "rb"),
                                    caption="Готово!")
            os.remove(data["name"] + '.pptx')
            logging.log(logging.INFO, slides)
            logging.log(logging.INFO, message.from_user.id)
            await state.finish()
        else:
            await PresentationStates.last_type.set()
            await bot.edit_message_caption(chat_id=message.chat.id, message_id=data["types_message"], caption="Тип №1",
                                           reply_markup=INLINE_MARKUP_TYPES)
            await bot.send_message(chat_id=message.chat.id, reply_to_message_id=data["types_message"],
                                   text="Следующий тип слайда?", reply_markup=ReplyKeyboardRemove())
            # await bot.send_photo(chat_id=message.chat.id, photo=open("types/0.png", "rb"), caption="Тип №1",
            #                      reply_markup=INLINE_MARKUP_TYPES)
    else:
        await message.answer(MESSAGES["make_slide"].format(now_number + 1))
        await PresentationStates.slides.set()


if __name__ == '__main__':
    logging.basicConfig(level=logging.INFO)
    start_polling(dispatcher=dp, skip_updates=False)
    # start_webhook(
    #     dispatcher=dp,
    #     webhook_path=WEBHOOK_PATH,
    #     skip_updates=False,
    #     on_startup=on_startup,
    #     on_shutdown=on_shutdown,
    #     host=WEBAPP_HOST,
    #     port=WEBAPP_PORT,
    # )

import pptx
from PIL import Image
prs = pptx.Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
slide.placeholders[0].text = "Текст"
pic = Image.open("11.jpeg")
pic = pic.resize((435, 330), Image.ANTIALIAS)
pic.save("a.jpeg")
slide.shapes.add_picture("a.jpeg", 1800000, 600000)
slide.placeholders[2].text = "Подтекст №2"
prs.save('_presentation.pptx')
import pptx
prs = pptx.Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.placeholders[0].text = "Текст"
prs.save("ggg.pptx")
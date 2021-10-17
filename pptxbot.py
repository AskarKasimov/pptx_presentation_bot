import ctypes
import logging
import pptx
import telebot
import os
import json
from config import SLIDE_TYPES, CONTENT, TOKEN, LOG_FILE, MARKUP

DESIGNS_NUMBER = len(os.listdir("designs")) - 1
TYPES_NUMBER = len(os.listdir("types"))
logging.basicConfig(filename=LOG_FILE, level=logging.INFO)
# inline клавиатура
INLINE_MARKUP_DESIGN = telebot.types.InlineKeyboardMarkup()
INLINE_MARKUP_DESIGN.row(
    telebot.types.InlineKeyboardButton(text='Назад', callback_data="back " + str(DESIGNS_NUMBER - 1) + " design"),
    telebot.types.InlineKeyboardButton(text='Вперед', callback_data="next 1" + " design"))
INLINE_MARKUP_DESIGN.add(telebot.types.InlineKeyboardButton(text='Этот', callback_data="0" + " none design"))



def content_help(message, typ):
    er = "Введите "
    if message.content_type == "photo":
        er = "Зачем мне картинка? Мне нужно получить "
    elif message.content_type == "document":
        er = "Давайте лучше без документов. Пришлите мне "
    elif message.content_type == "voice":
        er = "Очень красивый голос, но мне нужно прислать "
    elif message.content_type == "audio":
        er = "Потом послушаю... Сейчас мне бы лучше "
    elif message.content_type == "sticker":
        er = "Классный стикер, я себе тоже добавил... Но для презентации мне нужно отправить "
    elif message.content_type == "video" or message.content_type == "video_note":
        er = "Потом обязательно посмотрю, а сейчас вам нужно отправить мне "
    elif message.content_type == "location":
        er = "Были там? Красиво, наверное... Но я всё еще жду "
    elif message.content_type == "contact":
        er = "Ваш номерок? Обязательно созвонимся, но а сейчас мне требуется "
    return er + typ


def make_presentation(design, types, naming, title=None, subtitle=None):
    prs = pptx.Presentation(f"designs/{design}.pptx")
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.placeholders[0].text = str(title)
        slide.placeholders[1].text = str(subtitle) if subtitle else ""
    for txt in types:
        slide = prs.slides.add_slide(prs.slide_layouts[txt["type"]])
        for i in range(SLIDE_TYPES[txt["type"]]):
            slide.placeholders[i].text = str(txt[f"text{i}"])
    prs.save(naming)


class Bot:
    def __init__(self, bott):
        self._count_slides = 0
        self._count_text = 0
        self._help_type = 0

        self.prs_bot = bott
        self.prs_design = 0
        self.prs_title = ""
        self.prs_content = list()

    def start(self, message):
        print(str(message.from_user.username), str(message.text))
        self.prs_bot.send_message(message.chat.id,
                                  "Привет, я бот, помогающий сделать презентацию в привычном формате "
                                  ".pptx прямо в Telegram!")
        bot.send_message(message.chat.id,
                         "Каким будет дизайн вашей презентации?")
        bot.send_photo(message.chat.id, open("designs/demo/0.png", "rb"), f"Дизайн №1",
                       reply_markup=INLINE_MARKUP_DESIGN)

    def get_design(self, message):
        self.prs_design = message.text
        self.prs_bot.send_message(message.chat.id, "Сколько слайдов должно быть в вашей презентации?\n--\n"
                                                   "P.S. Указывайте количество слайдов, не учитывая титульный. "
                                                   "Заполнить титульный слайд можно "
                                                   "будет в конце создания презентации.")

        self.prs_bot.register_next_step_handler(message, self.get_count_slides)

    def get_count_slides(self, message):
        print(str(message.from_user.username), str(message.text))
        try:
            if int(message.text) > 0:
                self._count_slides = int(message.text)
                self.prs_bot.send_message(message.chat.id,
                                          "Сейчас нужно будет выбирать типы слайдов и "
                                          "вводить заголовки и тексты, добавлять картинки...\n"
                                          "Итак, тип первого слайда?")

                INLINE_MARKUP_TYPE = telebot.types.InlineKeyboardMarkup()
                INLINE_MARKUP_TYPE.row(
                    telebot.types.InlineKeyboardButton(text='Назад',
                                                       callback_data="back " + str(TYPES_NUMBER - 1) + " type " + str(self._count_slides) + " none"),
                    telebot.types.InlineKeyboardButton(text='Вперед', callback_data="next 1" + " type " + str(self._count_slides) + " none"))
                INLINE_MARKUP_TYPE.add(
                    telebot.types.InlineKeyboardButton(text='Этот', callback_data="0" + " none type " + str(self._count_slides) + " none"))
                bot.send_photo(message.chat.id, open("types/0.png", "rb"), f"Тип №1",
                               reply_markup=INLINE_MARKUP_TYPE)
            elif int(message.text) == 0:
                raise ValueError("Число слайдов не может быть равным нулю!")
            else:
                raise ValueError("Число слайдов не может быть отрицательным!")
        except ValueError as e:
            er = e if str(e)[0] != "i" else "Введите целое число!"
            print("Bot", er)
            self.prs_bot.reply_to(message, er)
            self.prs_bot.register_next_step_handler(message, self.get_count_slides)
        except TypeError:
            self.prs_bot.reply_to(message, content_help(message, "число"))
            print("Bot", content_help(message, "число"))
            self.prs_bot.register_next_step_handler(message, self.get_count_slides)

    def get_types(self, message):
        self._help_type = int(message.text.split()[0])
        self._count_slides = int(message.text.split()[1])
        if str(message.text.split()[2]) != "none":
            print(message.text.split()[2])
            try:
                self.prs_content = ctypes.cast(message.text.split()[2], ctypes.py_object).value
            except Exception:
                print("no!")
            print(ctypes.cast(message.text.split()[2], ctypes.py_object).value)
        else:
            self.prs_content = list()
        print("not lol")
        self._count_text = SLIDE_TYPES[self._help_type]
        self.prs_content.append({"type": self._help_type})
        self.prs_bot.send_message(message.chat.id, "Введите текст №1 слайда №1")

        self.prs_bot.register_next_step_handler(message, self.get_slides)

    def get_slides(self, message):
        print(str(message.from_user.username), str(message.text))
        if message.content_type != "text":
            self.prs_bot.reply_to(message, content_help(message, "текст"))
            print("Bot", content_help(message, "текст"))
            self.prs_bot.register_next_step_handler(message, self.get_slides)
        elif message.content_type == "text":
            index_help = SLIDE_TYPES[self._help_type] - self._count_text
            self.prs_content[-1][f"text{index_help}"] = message.text
            self._count_text -= 1
            if self._count_text == 0:
                self._count_slides -= 1
                if self._count_slides == 0:
                    self.prs_bot.send_message(message.chat.id, "Вам нужен титульный слайд?", reply_markup=MARKUP)
                    self.prs_bot.register_next_step_handler(message, self.title)
                    return
                self.prs_bot.send_message(message.chat.id, "Следующий слайд... Тип?")
                INLINE_MARKUP_TYPE = telebot.types.InlineKeyboardMarkup()
                INLINE_MARKUP_TYPE.row(
                    telebot.types.InlineKeyboardButton(text='Назад',
                                                       callback_data="back " + str(
                                                           TYPES_NUMBER - 1) + " type " + str(self._count_slides) + " " + str(id(self.prs_content))),
                    telebot.types.InlineKeyboardButton(text='Вперед', callback_data="next 1" + " type " + str(self._count_slides) + " " + str(id(self.prs_content))))
                INLINE_MARKUP_TYPE.add(
                    telebot.types.InlineKeyboardButton(text='Этот', callback_data="0" + " none type " + str(self._count_slides) + " " + str(id(self.prs_content))))
                bot.send_photo(message.chat.id, open("types/0.png", "rb"), f"Тип №1",
                               reply_markup=INLINE_MARKUP_TYPE)
                return
            self.prs_bot.send_message(message.chat.id,
                                      f"А теперь текст "
                                      f"№{SLIDE_TYPES[self._help_type] - self._count_text + 1} слайда "
                                      f"№{len(self.prs_content)}")
            print("Bot", "А теперь текст слайда")
            self.prs_bot.register_next_step_handler(message, self.get_slides)

    def title(self, message):
        if message.text == "Нет, не нужен":
            self.prs_bot.send_message(message.chat.id, "Приготовьтесь, сейчас вылетит презентация...",
                                      reply_markup=telebot.types.ReplyKeyboardRemove())
            logging.info(message.from_user.username + str(self.prs_content))
            make_presentation(self.prs_design, self.prs_content, message.from_user.username + '_presentation.pptx',
                              self.prs_title,
                              message.text)
            self.prs_bot.send_document(message.chat.id, open(message.from_user.username + '_presentation.pptx', 'rb'))
            os.remove(message.from_user.username + '_presentation.pptx')
        elif message.text == "Да, нужен":
            self.prs_bot.send_message(message.chat.id, "Введите текст заголовка титульного слайда",
                                      reply_markup=telebot.types.ReplyKeyboardRemove())
            self.prs_bot.register_next_step_handler(message, self.get_title)
        else:
            self.prs_bot.reply_to(message, "Выберете вариант ответа на экранчике под вашей клавиатурой")
            self.prs_bot.register_next_step_handler(message, self.title)

    def get_title(self, message):
        self.prs_title = message.text
        self.prs_bot.send_message(message.chat.id, "Вам нужен подзаголовок?", reply_markup=MARKUP)
        self.prs_bot.register_next_step_handler(message, self.get_subs)

    def get_subs(self, message):
        if message.text == "Нет, не нужен":
            self.prs_bot.send_message(message.chat.id, "Приготовьтесь, сейчас вылетит презентация...",
                                      reply_markup=telebot.types.ReplyKeyboardRemove())
            logging.info(message.from_user.username + str(self.prs_content))
            make_presentation(self.prs_design, self.prs_content, message.from_user.username + '_presentation.pptx',
                              self.prs_title)
            self.prs_bot.send_document(message.chat.id, open(message.from_user.username + '_presentation.pptx', 'rb'))
            os.remove(message.from_user.username + '_presentation.pptx')
        elif message.text == "Да, нужен":
            self.prs_bot.send_message(message.chat.id, "Введите текст подзаголовка титульного слайда",
                                      reply_markup=telebot.types.ReplyKeyboardRemove())
            self.prs_bot.register_next_step_handler(message, self.sending)
        else:
            self.prs_bot.reply_to(message, "Выберете вариант ответа на экранчике под вашей клавиатурой")
            self.prs_bot.register_next_step_handler(message, self.get_subs)

    def sending(self, message):
        self.prs_bot.send_message(message.chat.id, "Приготовьтесь, сейчас вылетит презентация...",
                                  reply_markup=telebot.types.ReplyKeyboardRemove())
        logging.info(message.from_user.username + str(self.prs_content))
        make_presentation(self.prs_design, self.prs_content, message.from_user.username + '_presentation.pptx',
                          self.prs_title,
                          message.text)
        self.prs_bot.send_document(message.chat.id, open(message.from_user.username + '_presentation.pptx', 'rb'))
        os.remove(message.from_user.username + '_presentation.pptx')


bot = telebot.TeleBot(TOKEN)


@bot.callback_query_handler(func=lambda call: call.data.split()[2] == "design")
def query_handler1(call):
    bot.answer_callback_query(call.id)
    if call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = telebot.types.InlineKeyboardMarkup()
        if call.data.split()[1] == str(DESIGNS_NUMBER - 1):
            but = telebot.types.InlineKeyboardButton(text='Вперед',
                                                     callback_data="next 0" + " design")
        else:
            but = telebot.types.InlineKeyboardButton(text='Вперед',
                                                     callback_data="next " +
                                                                   str(int(call.data.split()[1]) + 1) + " design")
        if call.data.split()[1] == "0":
            mark.row(telebot.types.InlineKeyboardButton(text='Назад',
                                                        callback_data="back " + str(DESIGNS_NUMBER - 1) + " design"),
                     but)
        else:
            mark.row(telebot.types.InlineKeyboardButton(text='Назад',
                                                        callback_data="back " + str(
                                                            int(call.data.split()[1]) - 1) + " design"),
                     but)
        mark.add(telebot.types.InlineKeyboardButton(text='Этот', callback_data=str(call.data.split()[1]) + " none design"))
        with open(f"designs/demo/{call.data.split()[1]}.png", "rb") as file:
            bot.edit_message_media(reply_markup=mark, chat_id=call.message.chat.id, message_id=call.message.message_id,
                                   media=telebot.types.InputMedia(type="photo", media=file))
            bot.edit_message_caption(f"Дизайн №{str(int(call.data.split()[1]) + 1)}", call.message.chat.id,
                                     call.message.message_id, reply_markup=mark)
    else:
        bot.edit_message_caption(f"Дизайн вашей презентации (№{str(int(call.data.split()[0]) + 1)}), выбранный вами",
                                 call.message.chat.id,
                                 call.message.message_id)
        call.message.text = int(call.data.split()[0])
        Bot(bot).get_design(call.message)


@bot.callback_query_handler(func=lambda call: call.data.split()[2] == "type")
def query_handler2(call):
    bot.answer_callback_query(call.id)
    if call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = telebot.types.InlineKeyboardMarkup()
        if call.data.split()[1] == str(TYPES_NUMBER - 1):
            but = telebot.types.InlineKeyboardButton(text='Вперед',
                                                     callback_data="next 0" + " type " + str(call.data.split()[3]) + " " + str(call.data.split()[4]))
        else:
            but = telebot.types.InlineKeyboardButton(text='Вперед',
                                                     callback_data="next " +
                                                                   str(int(call.data.split()[1]) + 1) + " type " + str(call.data.split()[3]) + " " + str(call.data.split()[4]))
        if call.data.split()[1] == "0":
            mark.row(telebot.types.InlineKeyboardButton(text='Назад',
                                                        callback_data="back " + str(TYPES_NUMBER - 1) + " type " + str(call.data.split()[3]) + " " + str(call.data.split()[4])),
                     but)
        else:
            mark.row(telebot.types.InlineKeyboardButton(text='Назад',
                                                        callback_data="back " + str(
                                                            int(call.data.split()[1]) - 1) + " type " + str(call.data.split()[3]) + " " + str(call.data.split()[4])),
                     but)
        mark.add(telebot.types.InlineKeyboardButton(text='Этот', callback_data=str(call.data.split()[1]) + " none type " + str(call.data.split()[3]) + " " + str(call.data.split()[4])))
        if call.data.split()[1] != "6":
            with open(f"types/{call.data.split()[1]}.png", "rb") as file:
                bot.edit_message_media(reply_markup=mark, chat_id=call.message.chat.id, message_id=call.message.message_id,
                                       media=telebot.types.InputMedia(type="photo", media=file))
                bot.edit_message_caption(f"Тип №{str(int(call.data.split()[1]) + 1)}", call.message.chat.id,
                                         call.message.message_id, reply_markup=mark)
        else:
            with open("types/7.png", "rb") as file:
                bot.edit_message_media(reply_markup=mark, chat_id=call.message.chat.id, message_id=call.message.message_id,
                                       media=telebot.types.InputMedia(type="photo", media=file))
                bot.edit_message_caption(f"Тип №{str(int(call.data.split()[1]) + 1)}", call.message.chat.id,
                                         call.message.message_id, reply_markup=mark)
    else:
        bot.edit_message_caption(f"Тип (№{str(int(call.data.split()[0]) + 1)}) слайда, выбранный вами",
                                 call.message.chat.id,
                                 call.message.message_id)
        print(call.data)
        call.message.text = str(call.data.split()[0]) + " " + str(call.data.split()[3]) + " " + str(call.data.split()[4])
        Bot(bot).get_types(call.message)


@bot.message_handler(commands=['start'])
def start(message):
    Bot(bot).start(message)


@bot.message_handler(content_types=CONTENT)
def index(message):
    bot.send_message(message.chat.id, "Нажмите /start чтобы начать")


bot.polling(none_stop=True)

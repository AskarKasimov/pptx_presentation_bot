import os
import pptx
import logging
from aiogram.utils.keyboard import InlineKeyboardBuilder
from aiogram import Router
from aiogram.filters import Command
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import StatesGroup, State
from aiogram.types import Message, ReplyKeyboardRemove, InlineKeyboardButton, CallbackQuery, InputMediaPhoto, ReplyKeyboardMarkup, KeyboardButton, ReplyKeyboardRemove, FSInputFile

MESSAGES = {
    "start": ["Привет, я бот, помогающий сделать презентацию в привычном формате .pptx прямо в Telegram!",
              "Введите /help, чтобы посмотреть список доступных команд"],
    "help": ["Вот список доступных команд:", "/new – Создать новую презентацию"],
    "presentation_start": "Введите название презентации (так будет называться .pptx файл)",
    "get_name": "Отлично! Ваш файл будет называться {}.pptx",
    "get_design_type": "Теперь нужно выбрать подходящий дизайн:",
    "get_slides_number": "Сколько слайдов должно быть в презентации?",
    "start_slides": "Итак, начнём заполнять слайды",
    "make_types": "Выберите тип слайда №{}",
    "make_slide": "Введите текст №{}",
    "make_pptx": "Принято! Отправил презентацию на кухню:)",
    "cancel": "Создание презентации отменено",
    "cancel_button": u'\U0000274C' + " Отменить создание презентации",

}

# prs = pptx.Presentation("designs/3.pptx")
#
# prs.slides.add_slide(prs.slide_layouts[8])
#
# shapes = prs.slides[0].shapes
#
# img = Image.open("designs/demo/2.png")
# img.thumbnail(size=(720, 390))
# buf = io.BytesIO()
# img.save(buf, format='PNG')
#
# image = shapes.add_picture(buf, left=0, top=0)
#
# image.left = (prs.slide_width - image.width) // 2
#
# prs.slides[0].placeholders[0].text = "gg"
#
# prs.save("GG.pptx")

ERRORS = {
    "int_required": "Ошибка. Введите натуральное число"
}

DESIGNS_NUMBER = len(os.listdir("designs")) - 1
SLIDE_TYPES = {
    0: 2,
    1: 2,
    2: 2,
    3: 3,
    4: 5,
    5: 1,
    6: 4,
    7: 3
}

TYPE_NUMBER = len(SLIDE_TYPES) - 1

INLINE_MARKUP_DESIGNS = InlineKeyboardBuilder()
INLINE_MARKUP_DESIGNS.row(InlineKeyboardButton(text="Назад", callback_data="back " + str(DESIGNS_NUMBER - 1)),
                          InlineKeyboardButton(text="Вперёд", callback_data="next 1"))
INLINE_MARKUP_DESIGNS.add(InlineKeyboardButton(text="Этот", callback_data="0"))
INLINE_MARKUP_DESIGNS.add(InlineKeyboardButton(
    text=MESSAGES["cancel_button"], callback_data="cancel"))

INLINE_MARKUP_TYPES = InlineKeyboardBuilder()
INLINE_MARKUP_TYPES.row(InlineKeyboardButton(text="Назад", callback_data="back " + str(TYPE_NUMBER)),
                        InlineKeyboardButton(text="Вперёд", callback_data="next 1"))
INLINE_MARKUP_TYPES.add(InlineKeyboardButton(text="Этот", callback_data="0"))
INLINE_MARKUP_TYPES.add(InlineKeyboardButton(
    text=MESSAGES["cancel_button"], callback_data="cancel"))

btns = [[KeyboardButton(text=MESSAGES["cancel_button"])]]
KEYBOARD_CANCEL = ReplyKeyboardMarkup(resize_keyboard=True, keyboard=btns)
# KEYBOARD_CANCEL.add(KeyboardButton(MESSAGES["cancel_button"]))

router = Router()


def make_presentation(design, slides, naming):
    prs = pptx.Presentation(f"designs/{design}.pptx")
    for slide_info in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[slide_info["type"]])
        for i in range(SLIDE_TYPES[slide_info["type"]]):
            slide.placeholders[i].text = str(slide_info["content"][i])
    prs.save(naming)


async def check_cancel(message: Message, state: FSMContext):
    """ Проверка кнопки клавиатуры на сброс """
    if message.text == MESSAGES["cancel_button"]:
        await message.answer(MESSAGES["cancel"], reply_markup=ReplyKeyboardRemove())
        await state.clear()
        return True
    return False


class PresentationStates(StatesGroup):
    name = State()
    design_type = State()
    slides_number = State()
    types_message = State()
    last_type = State()
    slides = State()


@router.message(Command(commands=["start"]))
async def bot_start(message: Message):
    """ Стандартная команда начала диалога """
    for msg in MESSAGES["start"]:
        await message.answer(msg)


@router.message(Command(commands=["help"]))
async def bot_help(message: Message):
    """ Показ пользователю доступных команд """
    for msg in MESSAGES["help"]:
        await message.answer(msg)


@router.message(Command(commands=["new"]))
async def presentation_start(message: Message, state: FSMContext):
    """ Начало создания презентации – запрос названия файла """
    try:
        logging.log(logging.INFO, str(message.from_user.id) +
                    " " + message.from_user.username)
    except Exception:
        logging.log(logging.INFO, message.from_user.id)
    await message.answer(MESSAGES["presentation_start"], reply_markup=KEYBOARD_CANCEL)
    await state.set_state(PresentationStates.name)


@router.message(PresentationStates.name)
async def get_filename(message: Message, state: FSMContext):
    """ Получение названия от пользователя """
    if await check_cancel(message, state):
        return
    await state.update_data(name=message.text)
    await message.answer(MESSAGES["get_name"].format(message.text))
    await state.set_state(PresentationStates.design_type)
    await message.answer(MESSAGES["get_design_type"], reply_markup=ReplyKeyboardRemove())
    await message.answer_photo(photo=FSInputFile("designs/demo/0.png"), caption="Дизайн №1",
                               reply_markup=INLINE_MARKUP_DESIGNS.as_markup())


@router.callback_query(PresentationStates.design_type)
async def inline_designs(call: CallbackQuery, state: FSMContext):
    """ Логика inline-клавиатуры под фото дизайнов """
    # await call.answer_callback_query(call.id)
    if call.data == "cancel":
        await call.message.edit_caption(caption=MESSAGES["cancel"],
                                        reply_markup=None)
        await state.clear()
    elif call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = InlineKeyboardBuilder()
        if call.data.split()[1] == str(DESIGNS_NUMBER - 1):
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next 0")
        else:
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next " + str(int(call.data.split()[1]) + 1))
        if call.data.split()[1] == "0":
            mark.row(InlineKeyboardButton(text='Назад', callback_data="back " + str(DESIGNS_NUMBER - 1)),
                     but)
        else:
            mark.row(InlineKeyboardButton(text='Назад',
                                          callback_data="back " + str(int(call.data.split()[1]) - 1)),
                     but)
        mark.add(InlineKeyboardButton(
            text='Этот', callback_data=call.data.split()[1]))
        mark.add(InlineKeyboardButton(
            text=MESSAGES["cancel_button"], callback_data="cancel"))
        # with open(f"designs/demo/{call.data.split()[1]}.png", "rb") as file:
        await call.message.edit_media(media=InputMediaPhoto(media=FSInputFile(f"designs/demo/{call.data.split()[1]}.png"), caption=f"Дизайн №{str(int(call.data.split()[1]) + 1)}"), reply_markup=mark.as_markup())
    else:
        await call.message.edit_caption(caption=f"Дизайн вашей презентации (№{str(int(call.data) + 1)}), выбранный вами")
        await state.update_data(design_type=int(call.data))
        await call.message.answer(MESSAGES["get_slides_number"], reply_markup=KEYBOARD_CANCEL)
        await state.set_state(PresentationStates.slides_number)


@router.message(PresentationStates.slides_number)
async def get_number(message: Message, state: FSMContext):
    """ Получение количества слайдов от пользователя """
    if await check_cancel(message, state):
        return
    try:
        if float(message.text) != float(int(message.text)):
            raise ValueError
        if int(message.text) <= 0:
            raise ValueError
        await state.update_data(slides_number=int(message.text))
        await message.answer(MESSAGES["start_slides"])
        await message.answer(MESSAGES["make_types"].format(1), reply_markup=ReplyKeyboardRemove())
        await state.update_data(types_message=0)
        await state.set_state(PresentationStates.last_type)
        await message.answer_photo(photo=FSInputFile("types/0.png"), caption="Тип №1",
                                   reply_markup=INLINE_MARKUP_TYPES.as_markup())
    except ValueError:
        await message.answer(ERRORS["int_required"])
        await state.set_state(PresentationStates.slides_number)


@router.callback_query(PresentationStates.last_type)
async def inline_types(call: CallbackQuery, state: FSMContext):
    """ Логика inline-клавиатуры под фото типов слайдов """
    # await bot.answer_callback_query(call.id)
    data = await state.get_data()
    if data["types_message"] == 0:
        await state.update_data(types_message=call.message.message_id)
    if call.data == "cancel":
        await call.message.edit_caption(caption=MESSAGES["cancel"],
                                        reply_markup=None)
        await state.clear()
    elif call.data.split()[0] == "next" or call.data.split()[0] == "back":
        mark = InlineKeyboardBuilder()
        if call.data.split()[1] == str(TYPE_NUMBER - 1):
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next 0")
        else:
            but = InlineKeyboardButton(text='Вперед',
                                       callback_data="next " + str(int(call.data.split()[1]) + 1))
        if call.data.split()[1] == "0":
            mark.row(InlineKeyboardButton(text='Назад', callback_data="back " + str(TYPE_NUMBER - 1)),
                     but)
        else:
            mark.row(InlineKeyboardButton(text='Назад',
                                          callback_data="back " + str(int(call.data.split()[1]) - 1)),
                     but)
        mark.add(InlineKeyboardButton(
            text='Этот', callback_data=call.data.split()[1]))
        mark.add(InlineKeyboardButton(
            text=MESSAGES["cancel_button"], callback_data="cancel"))
        try:
            # with open(f"types/{call.data.split()[1]}.png", "rb") as file:
            await call.message.edit_caption(media=InputMediaPhoto(media=FSInputFile(f"types/{call.data.split()[1]}.png"), caption=f"Тип №{str(int(call.data.split()[1]) + 1)}"), reply_markup=mark.as_markup())
        except FileNotFoundError:
            """ Так как нет 6-ого типа слайда """
            # with open(f"types/{str(int(call.data.split()[1]) + 1)}.png", "rb") as file:
            await call.message.edit_media(media=InputMediaPhoto(media=FSInputFile(f"types/{str(int(call.data.split()[1]) + 1)}.png"), caption=f"Тип №{str(int(call.data.split()[1]) + 1)}"), reply_markup=mark.as_markup())
    else:
        try:
            slides_len = len(data["slides"])
        except KeyError:
            slides_len = 0
        slides_len += 1
        await call.message.edit_caption(
            caption=f"Тип слайда (№{str(int(call.data) + 1)
                                    }), выбранный вами для слайда №{slides_len}",)
        await call.message.answer(MESSAGES["make_slide"].format(1), reply_markup=KEYBOARD_CANCEL)
        await state.update_data(last_type=int(call.data))
        await state.set_state(PresentationStates.slides)


@router.message(PresentationStates.slides)
async def make_slides(message: Message, state: FSMContext):
    if await check_cancel(message, state):
        return
    data = await state.get_data()
    edited = False
    try:
        slides = data["slides"]
        last_slide = slides[-1]
        if SLIDE_TYPES[last_slide["type"]] == len(last_slide["content"]):
            current_slide = {}
        else:
            current_slide = last_slide
            edited = True
    except KeyError:
        slides = []
        current_slide = {}
    current_slide["type"] = data["last_type"]

    if edited:
        current_slide["content"].append(message.text)
        slides[-1] = current_slide
    else:
        current_slide["content"] = [message.text]
        slides.append(current_slide)

    await state.update_data(slides=slides)

    needed_number = SLIDE_TYPES[slides[-1]["type"]]
    now_number = len(slides[-1]["content"])
    if needed_number == now_number:
        if len(slides) == data["slides_number"]:
            await message.answer(MESSAGES["make_pptx"], reply_markup=ReplyKeyboardRemove())
            make_presentation(data["design_type"],
                              slides, data["name"] + '.pptx')
            await message.answer_document(document=FSInputFile(data["name"] + '.pptx'),
                                          caption="Готово!")
            os.remove(data["name"] + '.pptx')
            logging.log(logging.INFO, slides)
            logging.log(logging.INFO, message.from_user.id)
            await state.clear()
        else:
            await message.answer("KAIFF")
            # await PresentationStates.last_type.set()
            # await bot.edit_message_caption(chat_id=message.chat.id, message_id=data["types_message"], caption="Тип №1",
            #                                reply_markup=INLINE_MARKUP_TYPES)
            # await bot.send_message(chat_id=message.chat.id, reply_to_message_id=data["types_message"],
            #                        text="Следующий тип слайда?", reply_markup=ReplyKeyboardRemove())
            # await bot.send_photo(chat_id=message.chat.id, photo=open("types/0.png", "rb"), caption="Тип №1",
            #                      reply_markup=INLINE_MARKUP_TYPES)
    else:
        await message.answer(MESSAGES["make_slide"].format(now_number + 1))
        await state.set_state(PresentationStates.slides)
